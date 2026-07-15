import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5) # 16:9 ratio
    
    # Theme colors (Sleek Dark Mode)
    BG_COLOR = RGBColor(15, 23, 42)      # Slate 900
    CARD_BG = RGBColor(30, 41, 59)       # Slate 800
    TEXT_LIGHT = RGBColor(241, 245, 249) # Slate 100
    TEXT_MUTED = RGBColor(148, 163, 184) # Slate 400
    ACCENT_INDIGO = RGBColor(99, 102, 241) # Indigo 500
    ACCENT_CYAN = RGBColor(34, 211, 238)  # Cyan 400
    BORDER_COLOR = RGBColor(51, 65, 85)   # Slate 700

    def apply_dark_background(slide):
        # Add a full screen rectangle for background
        bg_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
        )
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = BG_COLOR
        bg_shape.line.fill.background() # No border
        
    def add_title(slide, text):
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = TEXT_LIGHT
        p.alignment = PP_ALIGN.LEFT

    # --- Slide 1: Title Slide ---
    slide_layout = prs.slide_layouts[6] # Blank
    slide1 = prs.slides.add_slide(slide_layout)
    apply_dark_background(slide1)
    
    # Title Text Box
    title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.333), Inches(3.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "BookRAG"
    p.font.size = Pt(64)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN
    p.alignment = PP_ALIGN.LEFT
    
    p2 = tf.add_paragraph()
    p2.text = "Quality-First Agentic Retrieval-Augmented Generation"
    p2.font.size = Pt(28)
    p2.font.bold = True
    p2.font.color.rgb = TEXT_LIGHT
    p2.space_before = Pt(10)
    p2.alignment = PP_ALIGN.LEFT
    
    p3 = tf.add_paragraph()
    p3.text = "Grounded Q&A over Multimodal Corporate & Finance Corpus"
    p3.font.size = Pt(18)
    p3.font.color.rgb = TEXT_MUTED
    p3.space_before = Pt(15)
    
    p4 = tf.add_paragraph()
    p4.text = "Saad Faran  |  Presentation Day  |  July 15, 2026"
    p4.font.size = Pt(14)
    p4.font.color.rgb = ACCENT_INDIGO
    p4.space_before = Pt(40)

    # Helper for layout containers
    def add_card(slide, left, top, width, height, title="", border=True):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        if border:
            card.line.color.rgb = BORDER_COLOR
            card.line.width = Pt(1.5)
        else:
            card.line.fill.background()
        
        if title:
            tb = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.15), width - Inches(0.3), Inches(0.5))
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = ACCENT_CYAN
            return left + Inches(0.15), top + Inches(0.65), width - Inches(0.3), height - Inches(0.8)
        return left + Inches(0.15), top + Inches(0.15), width - Inches(0.3), height - Inches(0.3)

    # --- Slide 2: What is BookRAG? ---
    slide2 = prs.slides.add_slide(slide_layout)
    apply_dark_background(slide2)
    add_title(slide2, "What is BookRAG? (The Core \"What\")")
    
    # Left Column (Overview Card)
    cx, cy, cw, ch = add_card(slide2, Inches(0.5), Inches(1.3), Inches(6.0), Inches(5.6), "Project Overview")
    tb = slide2.shapes.add_textbox(cx, cy, cw, ch)
    tf = tb.text_frame
    tf.word_wrap = True
    
    points_left = [
        ("Grounded Answers with Zero Hallucination", 
         "Every response is verified by a strict fact-verification step and must trace back directly to source excerpts. The system rejects unverified claims."),
        ("Multimodal Corpus from Legal Sources", 
         "SEC filings (10-K/20-F), World Bank economic reports, arXiv papers, and public-domain books, with all tables and formats fully parsed."),
        ("Genuinely Agentic RAG System", 
         "Orchestrated via LangGraph with multi-stage query routing, hybrid RRF retrieval, factual evaluation, and automatic retry queries.")
    ]
    for i, (title, desc) in enumerate(points_left):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = "• " + title
        p.font.bold = True
        p.font.size = Pt(16)
        p.font.color.rgb = TEXT_LIGHT
        if i > 0: p.space_before = Pt(15)
        
        p_desc = tf.add_paragraph()
        p_desc.text = desc
        p_desc.font.size = Pt(13)
        p_desc.font.color.rgb = TEXT_MUTED
        p_desc.space_before = Pt(3)
        p_desc.level = 0
        
    # Right Column (Stats Cards)
    stats = [
        ("206+", "Documents", "SEC Filings, economic reports, and books."),
        ("40,000+", "Semantic Chunks", "Structure-aware segments parsed and embedded."),
        ("Free / Local", "Tech Stack", "Groq/ngrok inference, CPU-local embeddings, SQLite, ChromaDB.")
    ]
    for i, (val, label, note) in enumerate(stats):
        top_offset = Inches(1.3) + Inches(i * 1.9)
        scx, scy, scw, sch = add_card(slide2, Inches(6.8), top_offset, Inches(6.0), Inches(1.7), border=True)
        
        # Value
        tb_val = slide2.shapes.add_textbox(scx, scy, Inches(2.2), sch)
        tf_val = tb_val.text_frame
        tf_val.word_wrap = True
        p_val = tf_val.paragraphs[0]
        p_val.text = val
        p_val.font.size = Pt(36)
        p_val.font.bold = True
        p_val.font.color.rgb = ACCENT_CYAN
        
        # Label & Note
        tb_lbl = slide2.shapes.add_textbox(scx + Inches(2.3), scy, Inches(3.4), sch)
        tf_lbl = tb_lbl.text_frame
        tf_lbl.word_wrap = True
        p_lbl = tf_lbl.paragraphs[0]
        p_lbl.text = label
        p_lbl.font.size = Pt(18)
        p_lbl.font.bold = True
        p_lbl.font.color.rgb = TEXT_LIGHT
        
        p_note = tf_lbl.add_paragraph()
        p_note.text = note
        p_note.font.size = Pt(12)
        p_note.font.color.rgb = TEXT_MUTED
        p_note.space_before = Pt(5)

    # --- Slide 3: System Architecture ---
    slide3 = prs.slides.add_slide(slide_layout)
    apply_dark_background(slide3)
    add_title(slide3, "System Architecture: End-to-End Flow")
    
    # 3 columns for 3 stages
    stages = [
        ("1. Data Acquisition (acquire.py)", 
         "Downloads public and regulatory documents legally from SEC, World Bank, and arXiv. Keeps a JSON manifest to prevent duplicate downloads and support resume capabilities.",
         ["SEC EDGAR (HTML/Tables)", "World Bank Reports (PDF)", "arXiv q-fin Papers (PDF)"]),
        ("2. Multimodal Ingestion (ingest.py)", 
         "Parses text and financial tables (Markdown) with PyMuPDF, BeautifulSoup, and OCR fallbacks. Chunks text, embeds with BGE, and indexes to Chroma DB and BM25.",
         ["Structure-aware Chunking", "BGE Query/Passage Embeds", "Chroma (Dense) + BM25 (Sparse)"]),
        ("3. Query & UI (FastAPI & Next.js)", 
         "Processes user questions in real-time. LangGraph coordinates the routing, search, and validation nodes. SSE streams progress and trace logs to a rich Next.js UI.",
         ["FastAPI + SSE Stream Endpoint", "Multi-tenant SQLite database", "Live 3-Panel Next.js UI Dashboard"])
    ]
    for i, (title, desc, items) in enumerate(stages):
        col_x = Inches(0.5 + i * 4.2)
        cx, cy, cw, ch = add_card(slide3, col_x, Inches(1.3), Inches(3.9), Inches(5.6), title)
        
        tb = slide3.shapes.add_textbox(cx, cy, cw, ch)
        tf = tb.text_frame
        tf.word_wrap = True
        
        p_desc = tf.paragraphs[0]
        p_desc.text = desc
        p_desc.font.size = Pt(13)
        p_desc.font.color.rgb = TEXT_MUTED
        p_desc.space_after = Pt(20)
        
        for item in items:
            p_item = tf.add_paragraph()
            p_item.text = "✓ " + item
            p_item.font.size = Pt(14)
            p_item.font.bold = True
            p_item.font.color.rgb = TEXT_LIGHT
            p_item.space_before = Pt(8)

    # --- Slide 4: LangGraph Agentic Pipeline ---
    slide4 = prs.slides.add_slide(slide_layout)
    apply_dark_background(slide4)
    add_title(slide4, "LangGraph Agentic Pipeline: Zero-Hallucination Loop")
    
    # Draw a diagram box on left, details on right
    cx1, cy1, cw1, ch1 = add_card(slide4, Inches(0.5), Inches(1.3), Inches(5.5), Inches(5.6), "Pipeline Node Execution Flow")
    tb1 = slide4.shapes.add_textbox(cx1, cy1, cw1, ch1)
    tf1 = tb1.text_frame
    tf1.word_wrap = True
    
    diagram_nodes = [
        "1. rewrite_and_route: Processor model rewrites query.",
        "   └─ Route options: 'general' | 'tool' | 'search' | 'rag'",
        "2. retrieve: Performs hybrid vector + BM25 keyword retrieval.",
        "3. generate: Heavy LLM answers query using ONLY excerpts.",
        "4. evaluate_grounding (CRITICAL FACT CHECKER):",
        "   ├─ Verified -> cross_reference -> build_final_answer",
        "   └─ Unverified -> expand_query -> retrieve (loop back)",
        "5. retry_count guard: Caps loop iterations to exactly one retry."
    ]
    for i, line in enumerate(diagram_nodes):
        p = tf1.add_paragraph() if i > 0 else tf1.paragraphs[0]
        p.text = line
        p.font.size = Pt(13)
        if "evaluate_grounding" in line or "rewrite_and_route" in line:
            p.font.bold = True
            p.font.color.rgb = ACCENT_CYAN
        elif "Route options" in line or "Verified" in line or "Unverified" in line:
            p.font.color.rgb = TEXT_MUTED
        else:
            p.font.color.rgb = TEXT_LIGHT
        p.space_before = Pt(8)

    cx2, cy2, cw2, ch2 = add_card(slide4, Inches(6.3), Inches(1.3), Inches(6.5), Inches(5.6), "Key Mechanisms")
    tb2 = slide4.shapes.add_textbox(cx2, cy2, cw2, ch2)
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    
    mechanisms = [
        ("The Grounding Gate (evaluate_grounding)", 
         "An independent LLM call acts as a strict fact-checker, checking every factual claim. On failure, it forces a retry. If it fails again, it appends an honest disclaimer warning the user."),
        ("Query Expansion & Retry (expand_query)", 
         "If retrieval misses, the query expander generates a search string adding 4-6 synonyms or related financial terms to search broader files."),
        ("Consensus cross-referencing (cross_reference)", 
         "Compiles agreements and contradictions among retrieved sources to alert the user of differences in filing figures.")
    ]
    for i, (title, desc) in enumerate(mechanisms):
        p = tf2.add_paragraph() if i > 0 else tf2.paragraphs[0]
        p.text = "• " + title
        p.font.bold = True
        p.font.size = Pt(15)
        p.font.color.rgb = ACCENT_INDIGO
        if i > 0: p.space_before = Pt(15)
        
        p_desc = tf2.add_paragraph()
        p_desc.text = desc
        p_desc.font.size = Pt(13)
        p_desc.font.color.rgb = TEXT_MUTED
        p_desc.space_before = Pt(3)

    # --- Slide 5: Hybrid Retrieval & RRF ---
    slide5 = prs.slides.add_slide(slide_layout)
    apply_dark_background(slide5)
    add_title(slide5, "Hybrid Retrieval & Reciprocal Rank Fusion (RRF)")
    
    cx, cy, cw, ch = add_card(slide5, Inches(0.5), Inches(1.3), Inches(12.333), Inches(5.6), "Retrieval Logic")
    tb = slide5.shapes.add_textbox(cx, cy, cw, ch)
    tf = tb.text_frame
    tf.word_wrap = True
    
    # RRF Formula block
    p_f = tf.paragraphs[0]
    p_f.text = "Reciprocal Rank Fusion (RRF) Formula:"
    p_f.font.size = Pt(18)
    p_f.font.bold = True
    p_f.font.color.rgb = ACCENT_CYAN
    p_f.space_after = Pt(8)
    
    p_form = tf.add_paragraph()
    p_form.text = "RRF_Score(d) = sum_{m in M} 1 / (k + r_m(d))"
    p_form.font.size = Pt(22)
    p_form.font.bold = True
    p_form.font.color.rgb = TEXT_LIGHT
    p_form.alignment = PP_ALIGN.CENTER
    p_form.space_after = Pt(15)
    
    methods = [
        ("Dense Retrieval (ChromaDB Vector Store)", 
         "Queries Chroma with the BGE-embedded query using query-specific prefixes. Cosine similarity recovers semantic matches and conceptual context (top 15 candidates)."),
        ("Sparse Retrieval (BM25 Keyword Search)", 
         "Tokenizes and scores the query against the document corpus using BM25Okapi. Matches exact codes, tickers (e.g. 'NVDA'), and numeric tables (top 15 candidates)."),
        ("Rank Fusion (k = 60)", 
         "Fuses the rankings without score normalization issues. The top 5 fused chunks are sent to the LLM. On query retry, first-attempt candidate documents are merged and re-fused, ensuring a cumulative context pool.")
    ]
    for i, (title, desc) in enumerate(methods):
        p = tf.add_paragraph()
        p.text = "• " + title
        p.font.bold = True
        p.font.size = Pt(15)
        p.font.color.rgb = ACCENT_INDIGO
        p.space_before = Pt(10)
        
        p_desc = tf.add_paragraph()
        p_desc.text = desc
        p_desc.font.size = Pt(13)
        p_desc.font.color.rgb = TEXT_MUTED
        p_desc.space_before = Pt(2)

    # --- Slide 6: Multimodal Ingestion Pipeline ---
    slide6 = prs.slides.add_slide(slide_layout)
    apply_dark_background(slide6)
    add_title(slide6, "Multimodal Ingestion: Ingesting Tables & Scans")
    
    # 4 grid blocks
    grids = [
        ("PDF Ingestion (PyMuPDF / pdfplumber)", 
         "Extracts text per page. Incorporates native table extraction. Converts tables into Markdown format so numbers survive embedding and RAG constraints.", 
         Inches(0.5), Inches(1.3)),
        ("SEC HTML Ingestion (BeautifulSoup & lxml)", 
         "Parses XHTML filings. Removes structural layout tags, cleans spacer columns/rows, filters data tables, and isolates narrative text to prevent duplication.", 
         Inches(6.8), Inches(1.3)),
        ("OCR Fallback (pytesseract / rapidocr)", 
         "If a PDF page has <40 characters, it is classified as a scanned document. Converts page to image and runs OCR automatically to extract text.", 
         Inches(0.5), Inches(4.3)),
        ("Audio Transcription (faster-whisper)", 
         "Transcribes uploaded audio files (.mp3, .wav, .m4a) using CPU-optimized Whisper base. Merges transcripts into searchable project document chunks.", 
         Inches(6.8), Inches(4.3))
    ]
    for title, desc, col_x, row_y in grids:
        cx, cy, cw, ch = add_card(slide6, col_x, row_y, Inches(6.0), Inches(2.6), title)
        tb = slide6.shapes.add_textbox(cx, cy, cw, ch)
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT_MUTED

    # --- Slide 7: Database Design ---
    slide7 = prs.slides.add_slide(slide_layout)
    apply_dark_background(slide7)
    add_title(slide7, "Database Design & Observability Logs")
    
    # 2 columns: Schema and Logging
    cx1, cy1, cw1, ch1 = add_card(slide7, Inches(0.5), Inches(1.3), Inches(6.0), Inches(5.6), "SQLAlchemy SQLite Schema")
    tb1 = slide7.shapes.add_textbox(cx1, cy1, cw1, ch1)
    tf1 = tb1.text_frame
    tf1.word_wrap = True
    
    schema_tables = [
        ("users", "id (PK), email (unique), name, password_hash, profile (Text)"),
        ("auth_sessions", "id (PK), user_id (FK), refresh_hash, expires_at, revoked, user_agent"),
        ("projects", "id (PK), name, description, user_id (FK), created_at, updated_at"),
        ("project_files", "id (PK), project_id (FK), filename, path, kind, status, chunks, error"),
        ("chats", "id (PK), user_id (FK), project_id, title, summary, unsummarized_start"),
        ("messages", "id (PK), chat_id (FK), role, content, sources (JSON), trace (JSON)"),
        ("event_logs", "id (PK), correlation_id, step, payload (JSON), latency_ms, level, ts")
    ]
    for i, (table, cols) in enumerate(schema_tables):
        p = tf1.add_paragraph() if i > 0 else tf1.paragraphs[0]
        p.text = f"• {table}: "
        p.font.bold = True
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT_LIGHT
        if i > 0: p.space_before = Pt(6)
        
        run = p.add_run()
        run.text = cols
        run.font.size = Pt(12)
        run.font.color.rgb = TEXT_MUTED

    cx2, cy2, cw2, ch2 = add_card(slide7, Inches(6.8), Inches(1.3), Inches(6.0), Inches(5.6), "Observability & Step-by-Step Logging")
    tb2 = slide7.shapes.add_textbox(cx2, cy2, cw2, ch2)
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    
    logs_info = [
        ("Step-by-Step Logging", 
         "Every LangGraph node execution (node_start/node_end), API requests, errors, and tool calls are captured structured with a correlation_id."),
        ("Multi-Tenant Isolation", 
         "SQL data and chat logs are partitioned using the user_id context, ensuring user session security and clean tenant division."),
        ("Observability Analytics", 
         "Calculates execution metrics (routing ratios, tool usage, average processing speed per node, grounding verdicts) to expose on the web dashboard.")
    ]
    for i, (title, desc) in enumerate(logs_info):
        p = tf2.add_paragraph() if i > 0 else tf2.paragraphs[0]
        p.text = "• " + title
        p.font.bold = True
        p.font.size = Pt(15)
        p.font.color.rgb = ACCENT_CYAN
        if i > 0: p.space_before = Pt(12)
        
        p_desc = tf2.add_paragraph()
        p_desc.text = desc
        p_desc.font.size = Pt(13)
        p_desc.font.color.rgb = TEXT_MUTED
        p_desc.space_before = Pt(3)

    # --- Slide 8: Long-Term Memory & Context Sharing ---
    slide8 = prs.slides.add_slide(slide_layout)
    apply_dark_background(slide8)
    add_title(slide8, "Long-Term Memory & Project Context Sharing")
    
    cx, cy, cw, ch = add_card(slide8, Inches(0.5), Inches(1.3), Inches(12.333), Inches(5.6), "Memory & Context Sharing Layers")
    tb = slide5.shapes.add_textbox(cx, cy, cw, ch)
    tf = tb.text_frame
    tf.word_wrap = True
    
    memories = [
        ("Memory Layer 1: Rolling Chat Summary", 
         "Chat history can grow infinitely, causing token overload and slow LLM speeds. When unsummarized turns exceed 10, the router model compresses them into a 2-4 sentence summary. System context is built as: [ learned profile ] + [ rolling summary ] + [ 8 recent raw turns ]."),
        ("Memory Layer 2: Cross-Chat User Profile", 
         "Every 3 user turns, the model analyzes the conversation history to update a concise list of 6 durable facts about the user (goals, risk appetite, interests). Placed at the very top of context for personalization."),
        ("Context Sharing & Citation", 
         "Workspace files uploaded to a Project are embedded into a project-specific Chroma collection. Any chat created under that Project prioritises project document chunks during retrieval. Frontend renders source metadata (filename, page, RRF score) as interactive citation cards that trigger file retrieval.")
    ]
    for i, (title, desc) in enumerate(memories):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = "• " + title
        p.font.bold = True
        p.font.size = Pt(16)
        p.font.color.rgb = ACCENT_CYAN
        if i > 0: p.space_before = Pt(12)
        
        p_desc = tf.add_paragraph()
        p_desc.text = desc
        p_desc.font.size = Pt(13)
        p_desc.font.color.rgb = TEXT_MUTED
        p_desc.space_before = Pt(4)

    # --- Slide 9: Function Calling & Web Search ---
    slide9 = prs.slides.add_slide(slide_layout)
    apply_dark_background(slide9)
    add_title(slide9, "Function/Tool Calling & DuckDuckGo Search")
    
    cx1, cy1, cw1, ch1 = add_card(slide9, Inches(0.5), Inches(1.3), Inches(6.0), Inches(5.6), "Robust Tool Calling Registry")
    tb1 = slide9.shapes.add_textbox(cx1, cy1, cw1, ch1)
    tf1 = tb1.text_frame
    tf1.word_wrap = True
    
    tools_list = [
        ("Calculator", " ast-based evaluator pre-processed to convert human notations ('^' to '**', '15% of' to '* 0.01 *', 'x' to '*', commas in numbers)."),
        ("Weather", " current conditions via Open-Meteo with stop-word cleaning ('today', 'forecast') in geocoding inputs."),
        ("Currency Converter", " Frankfurter exchange rates mapped dynamically with a fuzzy lookup table for currency names/symbols (e.g. 'euros' → 'EUR')."),
        ("Stocks & Crypto", " Live stocks (Yahoo Finance) via browser-mimicking headers to bypass blocks, and crypto prices (CoinGecko)."),
    ]
    for i, (t_name, t_desc) in enumerate(tools_list):
        p = tf1.add_paragraph() if i > 0 else tf1.paragraphs[0]
        p.text = f"• {t_name}: "
        p.font.bold = True
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_LIGHT
        if i > 0: p.space_before = Pt(10)
        
        run = p.add_run()
        run.text = t_desc
        run.font.size = Pt(13)
        run.font.color.rgb = TEXT_MUTED

    cx2, cy2, cw2, ch2 = add_card(slide9, Inches(6.8), Inches(1.3), Inches(6.0), Inches(5.6), "Web Search Routing")
    tb2 = slide9.shapes.add_textbox(cx2, cy2, cw2, ch2)
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    
    web_info = [
        ("Live DuckDuckGo Search", 
         "If the query requires current real-world information (classified as 'search'), the pipeline queries DuckDuckGo dynamically for up to 5 results."),
        ("Automatic Citations", 
         "Synthesizes answers with bracketed inline citation annotations (e.g. [1], [2]). Coordinates with sources to attach URLs in the frontend UI."),
        ("Graceful Degradation", 
         "If the search package fails or network is offline, search failures degrade cleanly with direct instructions to the user.")
    ]
    for i, (title, desc) in enumerate(web_info):
        p = tf2.add_paragraph() if i > 0 else tf2.paragraphs[0]
        p.text = "• " + title
        p.font.bold = True
        p.font.size = Pt(15)
        p.font.color.rgb = ACCENT_INDIGO
        if i > 0: p.space_before = Pt(12)
        
        p_desc = tf2.add_paragraph()
        p_desc.text = desc
        p_desc.font.size = Pt(13)
        p_desc.font.color.rgb = TEXT_MUTED
        p_desc.space_before = Pt(3)

    # --- Slide 10: Model Context Protocol (MCP) ---
    slide10 = prs.slides.add_slide(slide_layout)
    apply_dark_background(slide10)
    add_title(slide10, "Model Context Protocol (MCP) Integration")
    
    cx, cy, cw, ch = add_card(slide10, Inches(0.5), Inches(1.3), Inches(12.333), Inches(5.6), "MCP Architecture")
    tb = slide10.shapes.add_textbox(cx, cy, cw, ch)
    tf = tb.text_frame
    tf.word_wrap = True
    
    mcp_details = [
        ("Standardized Tool Integration", 
         "Implements Anthropic's Model Context Protocol. MCP separates local data and functions from the LLM, creating a clean API boundaries between the model and client tools."),
        ("Dynamic Server Discovery", 
         "Reads standard 'mcp_config.json' on startup. Discovers tools served by external processes (e.g. finance_tools.py) over stdio using JSON-RPC, exposing tools like 'loan_payment', 'compound_interest', and 'rule_of_72'."),
        ("Stateless stdio sessions", 
         "Spawns stdio servers, executes tool calls, and cleans up processes. This stateless design prevents sub-process zombie leaks and daemon crashes during queries.")
    ]
    for i, (title, desc) in enumerate(mcp_details):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = "• " + title
        p.font.bold = True
        p.font.size = Pt(16)
        p.font.color.rgb = ACCENT_CYAN
        if i > 0: p.space_before = Pt(12)
        
        p_desc = tf.add_paragraph()
        p_desc.text = desc
        p_desc.font.size = Pt(13)
        p_desc.font.color.rgb = TEXT_MUTED
        p_desc.space_before = Pt(4)

    # --- Slide 11: Challenges & Solutions ---
    slide11 = prs.slides.add_slide(slide_layout)
    apply_dark_background(slide11)
    add_title(slide11, "Engineering Challenges & Solutions")
    
    cx1, cy1, cw1, ch1 = add_card(slide11, Inches(0.5), Inches(1.3), Inches(6.0), Inches(5.6), "System/LLM Challenges")
    tb1 = slide11.shapes.add_textbox(cx1, cy1, cw1, ch1)
    tf1 = tb1.text_frame
    tf1.word_wrap = True
    
    llm_challenges = [
        ("Groq Daily Token Caps & Limits", 
         "Groq free endpoints hit rate limits (429) under multi-slide execution. Fixed by integrating a self-hosted in-house GPU endpoint (OpenAI-compatible) running qwen3.5-2b."),
        ("Tool Argument Parsing Failures", 
         "Small 2B models generated currency names ('dollars') or temporal words ('Tokyo weather today') that crashed APIs. Solved with robust inputs maps and query cleaning in tools.py."),
        ("Oversized HTML Filings Ingestion", 
         "HTML table parsing swap-thrashed 8GB RAM machines. Fixed by enforcing a strict 250-chunk document cap and ignoring spacer layout tables.")
    ]
    for i, (title, desc) in enumerate(llm_challenges):
        p = tf1.add_paragraph() if i > 0 else tf1.paragraphs[0]
        p.text = "• " + title
        p.font.bold = True
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_LIGHT
        if i > 0: p.space_before = Pt(10)
        
        p_desc = tf1.add_paragraph()
        p_desc.text = desc
        p_desc.font.size = Pt(12.5)
        p_desc.font.color.rgb = TEXT_MUTED
        p_desc.space_before = Pt(2)

    cx2, cy2, cw2, ch2 = add_card(slide11, Inches(6.8), Inches(1.3), Inches(6.0), Inches(5.6), "RAG & UI Challenges")
    tb2 = slide11.shapes.add_textbox(cx2, cy2, cw2, ch2)
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    
    rag_challenges = [
        ("Missing Citations & Ignored Uploads", 
         "LLM classified files queries as 'general' and bypassed RAG. Fixed by updating rewrite_and_route rules, prioritising uploaded project files, and printing console tracebacks."),
        ("Unresponsive Live Flowchart", 
         "Next.js flowchart remained blank after launch. Solved by implementing uvicorn warmup hooks and adding automatic API polling until model loads."),
        ("Broken SSE Streaming protocol", 
         "Frontend split events on '\\n', but sse-starlette uses '\\r\\n', breaking stream parses. Solved with a single-line delimiter fix in App.jsx.")
    ]
    for i, (title, desc) in enumerate(rag_challenges):
        p = tf2.add_paragraph() if i > 0 else tf2.paragraphs[0]
        p.text = "• " + title
        p.font.bold = True
        p.font.size = Pt(14)
        p.font.color.rgb = ACCENT_CYAN
        if i > 0: p.space_before = Pt(10)
        
        p_desc = tf2.add_paragraph()
        p_desc.text = desc
        p_desc.font.size = Pt(12.5)
        p_desc.font.color.rgb = TEXT_MUTED
        p_desc.space_before = Pt(2)

    # --- Slide 12: Roadmap & Future ---
    slide12 = prs.slides.add_slide(slide_layout)
    apply_dark_background(slide12)
    add_title(slide12, "Project Roadmap & Future Enhancements")
    
    cx, cy, cw, ch = add_card(slide12, Inches(0.5), Inches(1.3), Inches(12.333), Inches(5.6), "Key Goals")
    tb = slide12.shapes.add_textbox(cx, cy, cw, ch)
    tf = tb.text_frame
    tf.word_wrap = True
    
    goals = [
        ("Layout-Aware Ingestion", "Integrate advanced PDF parsers (like Docling) to handle complex financial statements, footnotes, and multi-column tables cleanly."),
        ("Quantized Vector Search", "Incorporate product-quantization in ChromaDB for sub-millisecond retrieval speeds on datasets exceeding 1,000 documents."),
        ("Per-Token End-to-End Streaming", "Upgrade backend from node-level SSE streaming to full token-level streams for improved typing visual effects in the Next.js UI."),
        ("Evaluation Harness Integration", "Integrate quantitative evaluators (Ragas / TruLens) to continuously monitor retrieval hit-rate and grounding precision."),
        ("Hosted Cloud Deployment", "Deploy frontend to Vercel and dockerize uvicorn backend containers for scaling on AWS ECS.")
    ]
    for i, (title, desc) in enumerate(goals):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = "• " + title
        p.font.bold = True
        p.font.size = Pt(16)
        p.font.color.rgb = ACCENT_CYAN
        if i > 0: p.space_before = Pt(15)
        
        p_desc = tf.add_paragraph()
        p_desc.text = desc
        p_desc.font.size = Pt(13)
        p_desc.font.color.rgb = TEXT_MUTED
        p_desc.space_before = Pt(4)

    # Save
    prs.save("/Users/saadfaran/Downloads/BookRAG/BookRAG_Presentation.pptx")
    print("Presentation saved successfully at /Users/saadfaran/Downloads/BookRAG/BookRAG_Presentation.pptx")

if __name__ == "__main__":
    create_presentation()
